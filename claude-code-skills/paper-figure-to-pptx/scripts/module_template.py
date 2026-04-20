"""
Module Template — 复制此文件作为每个 Module 的起点。

核心模式：
  1) build_module(prs) 函数：往传入的 Presentation 里 add 一个 slide 并绘制
  2) __main__ 入口：单独运行时保存独立 pptx 便于审阅
  3) 不在 build_module 里调 prs.save() —— 保存权交给调用者
     （单独审阅时由 __main__ 保存，合并时由 assemble.py 保存）

这样每个 Module 的 .py 文件既能独立 python moduleN.py 出单页 pptx，
又能被 assemble.py 批量 import 合并到一个多页 pptx。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from helpers import (
    # slide 管理
    create_blank_slide,
    # 基础绘图
    add_tb, add_shape,
    add_arrow, add_bidir_arrow, add_line,
    add_formula_tb,
    # 复合元素
    add_hbrace, add_icon, add_token_row, add_dashed_note_box,
    # 配色
    C_BLUE, C_LIGHT_BLUE, C_PALE_BLUE,
    C_GREEN, C_LIGHT_GREEN, C_PALE_GREEN,
    C_ORANGE, C_LIGHT_ORANGE, C_PALE_ORANGE,
    C_PURPLE, C_LIGHT_PURPLE,
    C_GRAY_TEXT, C_GRAY_BORDER, C_GRAY_BG,
    C_RED, C_RED_BOLD, C_BLACK, C_WHITE,
    C_NAVY, C_BROWN, C_YELLOW,
)


# ============================================================
# 跨 Module 共用的坐标常量（合并时对齐的关键）
# 每个 Module 的 .py 都应该定义这套相同常量
# ============================================================
DRAW_X       = 0.3       # 绘图区左边距
DRAW_RIGHT   = 13.00     # 绘图区右边界
DRAW_W       = DRAW_RIGHT - DRAW_X

# 水平分区的关键 x 坐标（按 Fig2 示例，按你的图调整）
CASE_X       = 0.3       # 左侧 Case/子图标签区起点
CASE_W       = 1.30
TTFT_X       = 1.60      # 左区（如 TTFT Region）起点
RED_LINE_X   = 6.75      # 中间红虚线 / 关键分隔 x
TPOT_X       = 6.80      # 右区（如 TPOT Region）起点


def build_module(prs):
    """
    在传入的 Presentation 里新增一个 slide 并绘制 module 内容。

    调用方式：
        from pptx import Presentation
        prs = Presentation()
        build_module(prs)
        prs.save("output.pptx")  # 由调用者控制保存

    返回:
        新创建的 slide 对象
    """
    slide = create_blank_slide(prs)

    # ------------------------------------------------------------
    # Module 垂直布局（按阶段三比例估算填入）
    # ------------------------------------------------------------
    # 示例：横扁带 Module，占中段 3.0 英寸高
    MOD_Y = 1.8
    MOD_H = 3.0

    # ------------------------------------------------------------
    # 示例元素（删除后写你自己的内容）
    # ------------------------------------------------------------

    # 示例 1: Region 虚线外框 + 浅色填充
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
              TTFT_X, MOD_Y, RED_LINE_X - TTFT_X - 0.05, MOD_H,
              fill=C_PALE_BLUE, line_color=C_BLUE, line_w=1.0, dash=True)

    # 示例 2: Region 标题
    add_tb(slide, TTFT_X, MOD_Y + 0.08,
           RED_LINE_X - TTFT_X - 0.05, 0.30,
           "Region Title (Example)", fs=13, bold=True, color=C_BLUE)

    # 示例 3: Case 标签外置竖排（不用大卡片）
    add_tb(slide, CASE_X, MOD_Y + 0.80, CASE_W, 0.35,
           "Case X:", fs=15, bold=True, color=C_NAVY, align=PP_ALIGN.LEFT)
    add_formula_tb(slide, CASE_X, MOD_Y + 1.20, CASE_W, 0.40,
                   [("x", 'n'), ("u", 'sub'), (" = 0", 'n')],
                   fs=17, bold=True, italic=True,
                   color=C_BLACK, align=PP_ALIGN.LEFT)

    # 示例 4: 真上下标公式
    add_formula_tb(slide, TTFT_X + 0.2, MOD_Y + 2.0,
                   RED_LINE_X - TTFT_X - 0.4, 0.40,
                   [("TTFT", 'n'), ("u", 'sub'),
                    (" = ", 'n'),
                    ("T", 'n'), ("G2A", 'sup'), ("u→k", 'sub')],
                   fs=14, bold=True, italic=True, color=C_RED)

    # 示例 5: 图标占位符（用户后期替换为官方 SVG）
    add_icon(slide, CASE_X + 0.1, MOD_Y + 2.2, 0.55, 0.55, "[UAV]")

    # 示例 6: 水平大括号（手动构造，别用 LEFT_BRACE 旋转）
    add_hbrace(slide, TTFT_X + 0.3, RED_LINE_X - 0.1,
               MOD_Y + 1.6, color=C_RED, width=2.0)

    # 示例 7: 双向箭头（跨 token 周期）
    add_bidir_arrow(slide, TPOT_X + 0.5, MOD_Y + 2.5,
                    DRAW_RIGHT - 0.5, MOD_Y + 2.5,
                    color=C_RED, width=1.75)

    # ------------------------------------------------------------
    # 关键装饰线最后画（避免被 Region 填充遮挡）
    # ------------------------------------------------------------
    # 红色垂直虚线：精确对齐关键时间点 x 坐标
    add_line(slide, RED_LINE_X, MOD_Y, RED_LINE_X, MOD_Y + MOD_H + 0.3,
             color=C_RED, width=1.75, dash=True)

    # ------------------------------------------------------------
    # Module 左下角标签（方便合并后识别是哪一块）
    # ------------------------------------------------------------
    add_tb(slide, 0.30, 7.15, 6.0, 0.28,
           "Module X / 模块描述",
           fs=9, italic=True, color=C_GRAY_TEXT, align=PP_ALIGN.LEFT)

    return slide


if __name__ == '__main__':
    # 单独运行：生成一个单 slide pptx 用于审阅
    prs = Presentation()
    build_module(prs)
    out = "Module_Template.pptx"
    prs.save(out)
    print(f"Saved: {out}")
