"""
paper-figure-to-pptx helpers
===============================
学术论文架构图 → 可编辑 PPT 的完整绘图工具箱。

所有 helper 消化了 python-pptx 的底层 XML 细节，Module 代码里只需要
import 这些函数即可，不用再手动写 etree/qn 操作。

用法：
    from helpers import *

    slide = create_blank_slide(prs)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1, 2, 3, 1,
              fill=C_PALE_BLUE, line_color=C_BLUE, line_w=1.0, dash=True)
    add_formula_tb(slide, 1, 3, 3, 0.4,
                   [("T", 'n'), ("u", 'sub'), (" = ", 'n'),
                    ("x", 'n'), ("sat", 'sup')],
                   fs=14, color=C_RED)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ============================================================
# 配色常量
# ============================================================
# 主色系（论文架构图常见语义色）
C_BLUE         = RGBColor(0x2B, 0x6C, 0xB0)   # 上行/输入/Prefill
C_LIGHT_BLUE   = RGBColor(0xBE, 0xDB, 0xF8)
C_PALE_BLUE    = RGBColor(0xE7, 0xF0, 0xFA)   # Region 浅蓝填充

C_GREEN        = RGBColor(0x2F, 0x9E, 0x44)   # 下行/交付/A2G
C_LIGHT_GREEN  = RGBColor(0xB2, 0xE0, 0xC1)
C_PALE_GREEN   = RGBColor(0xE6, 0xF4, 0xEA)

C_ORANGE       = RGBColor(0xE0, 0x77, 0x16)   # 计算/Decode/Compressor
C_LIGHT_ORANGE = RGBColor(0xFB, 0xD3, 0x8B)
C_PALE_ORANGE  = RGBColor(0xFD, 0xEE, 0xD5)

C_PURPLE       = RGBColor(0x7E, 0x3A, 0xF2)   # 卫星/Cloud/A2S/S2A
C_LIGHT_PURPLE = RGBColor(0xD6, 0xBC, 0xFA)   # ISL↑ / ISL↓

# 中性色
C_GRAY_TEXT    = RGBColor(0x4B, 0x55, 0x63)   # 灰色文字/辅助线
C_GRAY_BORDER  = RGBColor(0x94, 0xA3, 0xB8)   # 虚线边框/灰色箭头
C_GRAY_BG      = RGBColor(0xF1, 0xF5, 0xF9)   # 灰色背景
C_GRAY_LIGHT   = RGBColor(0xE2, 0xE8, 0xF0)

# 强调色
C_RED          = RGBColor(0xC0, 0x1C, 0x1C)   # 关键公式/红虚线/大括号/★
C_RED_BOLD     = RGBColor(0xA3, 0x16, 0x16)   # Decode bottleneck 粗边框

# 基础色
C_BLACK        = RGBColor(0x1F, 0x29, 0x37)
C_WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
C_YELLOW       = RGBColor(0xFA, 0xCC, 0x15)   # 决策菱形

# Case 标签区分色
C_NAVY         = RGBColor(0x1E, 0x3A, 0x6E)   # Case A 主题色
C_BROWN        = RGBColor(0x8B, 0x45, 0x13)   # Case B 主题色


# ============================================================
# 内部工具（用户无需直接调用）
# ============================================================
def _set_margins(tf, m=0):
    """把 text_frame 的四边 margin 都设为 m（EMU）。默认 0 让文字紧贴边界。"""
    tf.margin_left = Emu(m)
    tf.margin_right = Emu(m)
    tf.margin_top = Emu(m)
    tf.margin_bottom = Emu(m)


def _apply_ends(conn, tail=True, head=False, dash=False):
    """在 connector 的 XML 上设置箭头端点（tailEnd/headEnd）和 dash 样式。"""
    ln = conn.line._get_or_add_ln()
    for t in ln.findall(qn('a:tailEnd')):
        ln.remove(t)
    for h in ln.findall(qn('a:headEnd')):
        ln.remove(h)
    if tail:
        te = etree.SubElement(ln, qn('a:tailEnd'))
        te.set('type', 'triangle')
        te.set('w', 'med')
        te.set('h', 'med')
    if head:
        he = etree.SubElement(ln, qn('a:headEnd'))
        he.set('type', 'triangle')
        he.set('w', 'med')
        he.set('h', 'med')
    if dash:
        for pd in ln.findall(qn('a:prstDash')):
            ln.remove(pd)
        prstDash = etree.SubElement(ln, qn('a:prstDash'))
        prstDash.set('val', 'dash')


# ============================================================
# 基础工具
# ============================================================
def create_blank_slide(prs):
    """
    在 Presentation 里新增一个空白 16:9 slide 并返回。
    若 prs 尚未设置尺寸，同时设为 13.333 × 7.5 英寸。
    """
    if prs.slide_width != Inches(13.333):
        prs.slide_width = Inches(13.333)
    if prs.slide_height != Inches(7.5):
        prs.slide_height = Inches(7.5)
    return prs.slides.add_slide(prs.slide_layouts[6])


# ============================================================
# 文本与形状
# ============================================================
def add_tb(slide, x, y, w, h, text, fs=10, bold=False, color=C_BLACK,
           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
           italic=False, font="Calibri"):
    """
    添加文本框。坐标单位为英寸（inches）。

    参数:
        slide: pptx slide 对象
        x, y, w, h: 位置和尺寸（英寸）
        text: 文字内容
        fs: 字号 pt
        bold, italic: 粗体/斜体
        color: RGBColor（可用 C_* 常量）
        align: PP_ALIGN.LEFT / CENTER / RIGHT
        anchor: MSO_ANCHOR.TOP / MIDDLE / BOTTOM
        font: 字体名
    """
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    _set_margins(tf)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fs)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return tb


def add_shape(slide, stype, x, y, w, h, fill=None, line_color=None,
              line_w=1.0, dash=False, no_line=False, text=None, fs=9,
              bold=False, italic=False, text_color=C_BLACK):
    """
    添加形状（矩形、圆角矩形、椭圆、菱形等）。

    参数:
        stype: MSO_SHAPE.RECTANGLE / ROUNDED_RECTANGLE / OVAL / DIAMOND 等
        x, y, w, h: 位置和尺寸（英寸）
        fill: 填充色。None 表示无填充（透明）
        line_color: 边框色。None 表示默认黑
        line_w: 边框粗细 pt
        dash: True 边框为虚线
        no_line: True 无边框
        text: 可选文字。如果给定，同时在形状内部加居中文字
        fs, bold, italic, text_color: 内部文字格式
    """
    s = slide.shapes.add_shape(stype, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if no_line:
        s.line.fill.background()
    else:
        if line_color is not None:
            s.line.color.rgb = line_color
        s.line.width = Pt(line_w)
        if dash:
            ln = s.line._get_or_add_ln()
            for pd in ln.findall(qn('a:prstDash')):
                ln.remove(pd)
            prstDash = etree.SubElement(ln, qn('a:prstDash'))
            prstDash.set('val', 'dash')
    if text is not None:
        tf = s.text_frame
        _set_margins(tf, 18000)
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(fs)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = text_color
        run.font.name = "Calibri"
    return s


# ============================================================
# 线条与箭头
# ============================================================
def add_arrow(slide, x1, y1, x2, y2, color=C_BLACK, width=1.25, dash=False):
    """
    从 (x1,y1) 到 (x2,y2) 的单向箭头。箭头在 (x2,y2) 端。

    坐标单位英寸。dash=True 画虚线箭头。
    """
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    _apply_ends(conn, tail=True, head=False, dash=dash)
    return conn


def add_bidir_arrow(slide, x1, y1, x2, y2, color=C_RED, width=1.75):
    """双向箭头（两端都有箭头符号）。用于表示周期/范围等对称语义。"""
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    _apply_ends(conn, tail=True, head=True, dash=False)
    return conn


def add_line(slide, x1, y1, x2, y2, color=C_BLACK, width=1.0, dash=False):
    """
    无箭头直线。

    典型用途：
        - 红色垂直虚线（作为时间点标记）
        - 水平大括号的组成线
        - 分隔线 / 参考线
    """
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    _apply_ends(conn, tail=False, head=False, dash=dash)
    return conn


# ============================================================
# 公式（真上下标）
# ============================================================
def add_formula_tb(slide, x, y, w, h, parts, fs=14, color=C_BLACK,
                   bold=True, italic=True, align=PP_ALIGN.CENTER):
    """
    带真正上下标的公式文本框。

    PPT 打开后 sup/sub 是原生上下标字符，可直接用字体对话框或公式编辑器修改。

    参数:
        parts: list of (text, style) 元组
               style ∈ {'n' (普通), 'sup' (上标), 'sub' (下标)}

    示例:
        # TTFT_u = T^{G2A}_{u→k}
        parts = [
            ("TTFT", 'n'), ("u", 'sub'),
            (" = ", 'n'),
            ("T", 'n'), ("G2A", 'sup'), ("u→k", 'sub'),
        ]
        add_formula_tb(slide, 1, 2, 5, 0.4, parts, fs=14, color=C_RED)

    注意事项:
        - 连续的下标内容应合并成单个 run: ("s0→su", 'sub') 而非
          ("s", 'sub'), ("0", 'sub'), ("→s", 'sub'), ("u", 'sub')
          否则字符间距会异常
        - Unicode 下标字符（如 ₀, ᵤ）不要和 baseline 上下标混用
        - baseline 不是完美数学排版（n 同时带上下标时不会垂直对齐），
          关键公式建议在 PPT 里用公式编辑器重写
    """
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    _set_margins(tf)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    for text, style in parts:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(fs)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = "Calibri"
        rPr = run._r.get_or_add_rPr()
        if style == 'sup':
            rPr.set('baseline', '30000')
        elif style == 'sub':
            rPr.set('baseline', '-25000')
    return tb


# ============================================================
# 复合元素
# ============================================================
def add_hbrace(slide, x1, x2, y, color=C_RED, width=2.0, tip_depth=0.12):
    """
    水平大括号（手动 4 线构造，避免 LEFT_BRACE 旋转的坐标失控问题）。

    结构：
        │───────┬───────│     ← 两端上翘 + 主水平线
                ▼             ← 中间向下尖

    用于括住时序条等水平元素，中间的下尖指向下方公式/说明。

    参数:
        x1, x2: 大括号左右端点 x 坐标
        y: 大括号主水平线的 y 坐标
        tip_depth: 中间下尖的长度
    """
    # 主水平线
    add_line(slide, x1, y, x2, y, color=color, width=width)
    # 两端向上的短竖线
    add_line(slide, x1, y, x1, y - 0.08, color=color, width=width)
    add_line(slide, x2, y, x2, y - 0.08, color=color, width=width)
    # 中间向下的短竖线（尖端）
    mid_x = (x1 + x2) / 2
    add_line(slide, mid_x, y, mid_x, y + tip_depth, color=color, width=width)


def add_icon(slide, x, y, w, h, label, border_color=C_GRAY_TEXT, border_w=2.0):
    """
    图标占位符：加粗边框的圆角矩形 + [Type] 标签。

    用于无法用 python-pptx 画出来的复杂图标（UAV、卫星、手机、
    神经网络节点、数据库图标等）。用户后期在 PPT 里用官方 SVG 图标库
    （Cisco/AWS/Azure/Google Cloud/BioRender）替换这些占位符。

    推荐 label 格式：`[UAV]`, `[Phone]`, `[Satellite]`, `[Text]`, `[Image]`
    统一方括号 + 类型名，便于用户后期批量 find-and-replace。
    """
    return add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h,
        fill=C_WHITE, line_color=border_color, line_w=border_w,
        text=label, fs=9, bold=True, text_color=border_color
    )


def add_token_row(slide, x, y, count, unit_w=0.10, gap=0.02, h=0.35,
                  colors=None):
    """
    水平 token 方块序列。用于表示 LLM 输入/输出 token 流。

    参数:
        count: token 数量
        unit_w: 每个方块宽度
        gap: 方块之间间隔
        h: 方块高度
        colors: 单一 RGBColor（所有方块同色）或 list（按索引循环）

    返回:
        (shapes, total_width) — shapes 是创建的方块列表，total_width
        是整个 token 序列的总宽度（便于后续定位右边界）
    """
    if colors is None:
        colors = C_BLUE
    shapes = []
    for i in range(count):
        c = colors[i % len(colors)] if isinstance(colors, list) else colors
        s = add_shape(
            slide, MSO_SHAPE.RECTANGLE,
            x + i * (unit_w + gap), y, unit_w, h,
            fill=c, line_color=C_WHITE, line_w=0.5
        )
        shapes.append(s)
    total_w = count * unit_w + (count - 1) * gap
    return shapes, total_w


def add_dashed_note_box(slide, x, y, w, h, text, fs=8,
                        bg=C_GRAY_BG, border=C_GRAY_BORDER,
                        text_color=C_GRAY_TEXT):
    """
    灰色虚线说明框。用于图中辅助说明（如 "↓η : ↓Latency but ↓Quality"、
    "A2S → ISL → Satellite s_u" 这类小字注释）。

    视觉：浅灰填充 + 灰色虚线边框 + 灰色斜体文字。
    """
    box = add_shape(
        slide, MSO_SHAPE.RECTANGLE, x, y, w, h,
        fill=bg, line_color=border, line_w=0.75, dash=True
    )
    add_tb(
        slide, x + 0.05, y + 0.03, w - 0.10, h - 0.06,
        text, fs=fs, color=text_color, italic=True
    )
    return box
